# -*- coding: utf-8 -*-
"""
Created on Tue Oct 16 14:05:50 2018

@author: peeraponge
"""
#Workflow(Slide User guide)
#Include table and plot in presentation
#

#%% Import Libraries
from sklearn.preprocessing import PolynomialFeatures
from sklearn.linear_model import LinearRegression
from sklearn.metrics import r2_score
import pandas as pd
import numpy as np
import itertools
import os
from decimal import Decimal

#%% Import for plotting
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D
from matplotlib import cm
from pandas.tools.plotting import table

#%%Import for PPT
from pptx import Presentation # สำหรับ สร้าง ppt
from pptx.util import Inches # สำหรับจัดตำแหน่งของ ppt
from pptx.dml.color import RGBColor # สำหรับ จัดสีของ ppt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

#%% Function: Generate polynomial equation
def Poly_equation(feature, coef, intercept, power):
    '''
    Return fitted polynomial equation as a string
    
    feature: list of feature in dataset
    power: degree of each feature in each term (can get from poly.powers_)
    '''
    poly_string = ""
    
    for i in range(len(coef)): # create polynomial term
        
        #Coefficients
        if i == 0:
            term_string = "y = %.2E" % Decimal(coef[i])
        elif coef[i] >= 0: # add + sign in front of coef
            term_string = "+%.2E" % Decimal(coef[i])
        else:
            term_string = "%.2E" % Decimal(coef[i])
        
        #Powers
        feature_order = 0
        for power_iter in power[i]: # power for each feature
            if power_iter == 1 : #degree of that feature = 1
                term_string += '*' + str(feature[feature_order])
            elif power_iter > 1 : #degree of that feature > 1
                term_string += '*' + str(feature[feature_order]) + '^' + str(power_iter)
            feature_order += 1
        poly_string += term_string
    
    #Add intercept
    if intercept >= 0:
        poly_string += "+"
    poly_string += "%.2E" % Decimal(intercept)
    
    return poly_string

#%%Function: Generate log equation
def Log_equation(feature,a,b):
    # Desired equation: y = bln(x)+a
    log_string = "y = %.2E*ln(%s)" % (Decimal(b),feature)
    if b >= 0:
        log_string += "+"
    log_string += "%.2E" % Decimal(a)
    return log_string

#%% Function: Generate exponential equation
def Exp_equation(feature,d,e):
    # Desired equation: y = exp(e)*exp(d*x)
    exp_string = "y = %.2E*exp(%.2E*%s)" % (Decimal(np.exp(e)),Decimal(d),feature)
    return exp_string

#%% Function: Generate power equation
def Power_equation(feature,f,g):
    # Desired equation: y = exp(g)*x^f
    pow_string = "y = %.2E*%s^(%.2E)" % (Decimal(np.exp(g)),feature,Decimal(f))
    return pow_string

#%% Function: generate 2D plot and calculate R-squared
def R2_plot_2d(xg,y,nPoly,m,c,figname,feature,y_label):
    '''
    This function generate logarithm/exponential/power 2D plot, and save it into .png
    Also, return R-squared value
    '''
    fig = plt.figure(figsize=(5.5,5))
    ax = fig.add_subplot(111)
    ax.scatter(xg,y) #Scatter plot
    
    xx = np.linspace(min(xg),max(xg),200)
    try:
        if nPoly == 0:
            yy = c+m*np.log(xx)
            ypred = c+m*np.log(xg)
        elif nPoly == -1:
            yy = np.exp(c)*np.exp(m*xx)
            ypred = np.exp(c)*np.exp(m*xg)
        else:
            yy = np.exp(c)*(xx**m)
            ypred = np.exp(c)*(xg**m)
        ax.plot(xx,yy,color = "black") #Line plot
        r2 = "%.3f" %r2_score(y,ypred) #Calculate R-squared
    except ValueError:
        print("Warning: Cannot find equation for %s,%s" %(nPoly,feature))
        r2 = "Error"
        
    ax.set_xlabel(feature)
    ax.set_ylabel(y_label,rotation = 0)
    ax.xaxis.set_label_coords(1.05, -0.025)
    ax.yaxis.set_label_coords(-0.025, 1.05)
    
    fig.savefig(figname, dpi=200) #Save figure as .png
        
    return r2

#%% Function: Add equation, R-squared and plot into slide presentation
def Add_slide(num_poly,feature,img_path,equ,r2,filename):

    sld1 = prs.slides.add_slide(prs.slide_layouts[5])
    #Add title
    if num_poly > 0:
        sld1.placeholders[12].text = "%s, Polynomial degree %d" %(feature, num_poly)
    elif num_poly == 0:
        sld1.placeholders[12].text = "%s, Logarithm" %(feature)
    elif num_poly == -1:
        sld1.placeholders[12].text = "%s, Exponential" %(feature)
    else:
        sld1.placeholders[12].text = "%s, Power" %(feature)
    #Add text box for equation and R-squared
    txBox = sld1.shapes.add_textbox(left=Inches(5.5), top=Inches(1.75), width=Inches(4), height=Inches(4))
    tf = txBox.text_frame
    tf.text = "Equation:\n%s\n" % equ #Add equation
    p = tf.add_paragraph()
    p.text = "R-squared: %s\n" %r2 #Add R-squared
    tf.word_wrap = True
    
    #Add plot
    if img_path == "none":
        #Add textbox to display "Cannot display graph"
        tx2 = sld1.shapes.add_textbox(left = Inches(0.5), top = Inches(1.75), width=Inches(5),height=Inches(4.75))
        tf2 = tx2.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf2.paragraphs[0]
        p.text = "Cannot display graph"
        p.alignment = PP_ALIGN.CENTER
        tx2.line.color.rgb = RGBColor(0,0,0)
    else:
        sld1.shapes.add_picture(img_path, left = Inches(0), top = Inches(1.75), width=Inches(5.5))
    
    #Delete image file
    if os.path.exists(img_path):
        os.remove(img_path) 

#%%

# Main Code Start here

#%% Load data from excel and initiate variables
        
excelFile         = "Input.xlsx"
excelSheet        = "Input"
SlideTemplateName = "Template.pptx"

#usecols = int(input("Number of independent variables: "))
deg = 3 #Maximum polynomial degree
df = pd.read_excel(io         = excelFile,
                   sheet_name = excelSheet,
                   header     = 0,
#                   usecols     = usecols
                  )
df_X = df.drop(df.columns[[0]], axis=1) #X input
y = df[df.columns[0]] #y input
y_label = df.columns[0]
prs = Presentation(SlideTemplateName) #Create pptx presentation
record = pd.DataFrame(columns=['num_feature','feature','degree','Rsquare']) #Summary record
n = df_X.shape[1]

#%% Train polynomial regression model
for num_feature in range(1,n+1):
    print("\n start running number of feature:"+str(num_feature))
    feature_list = list(itertools.combinations(df_X.columns,num_feature)) # get all combination of features
    feature = [None] * num_feature
    
    for feature in feature_list:
        print("\n features:"+str(feature))
        X = df_X[list(feature)]
        R_square = []
        
        #Perform regression from polynomial degree 1 to deg
        for nPoly in range(1, deg+1):
            poly = PolynomialFeatures(degree=nPoly, include_bias=None)
            # transform features to polynomial terms
            X_ = poly.fit_transform(X)  

            # train linear regression model
            lg = LinearRegression()
            lg.fit(X_, y)
            
            # Acquire equation and print R-squared score
            coef = Poly_equation(feature, lg.coef_, lg.intercept_, poly.powers_)
            print("Degree %d, R-square= %s" %(nPoly,str(lg.score(X_,y))))
            
            #%%Plotting polynomial graph
            if num_feature == 1:
                #2D plot for feature = 1
                fig = plt.figure(figsize=(5.5,5))
                ax = fig.add_subplot(111)
                xg = df_X.loc[:,feature[0]]
                ax.scatter(xg,y) #Scatter plot
                xx = np.linspace(min(xg),max(xg),200)
                XX_ = poly.fit_transform(pd.DataFrame(xx))
                yy = lg.predict(XX_)
                ax.plot(xx,yy,color = "black") #Line plot
                
                ax.set_xlabel(feature[0])
                ax.set_ylabel(y_label, rotation = 0)
                ax.xaxis.set_label_coords(1.05, -0.025)
                ax.yaxis.set_label_coords(-0.025, 1.05)
                
                figname = "%s_poly%d_%s.png" %(excelSheet,nPoly,str(feature))
                fig.savefig(figname, dpi=200) #Save plot
                
            elif num_feature == 2:
                #3D surface and scatter plot for feature = 2
                fig = plt.figure(figsize=(5.5,5))
                ax = fig.add_subplot(111, projection='3d')
                xg = df_X.loc[:,feature[0]]
                yg = df_X.loc[:,feature[1]]
                ax.scatter(xg, yg, y, zdir='z', s=20) #scatter plot
                
                xx = np.linspace(min(xg),max(xg),200)
                yy = np.linspace(min(yg),max(yg),200)
                xxx, yyy = np.meshgrid(xx, yy)
                XX = np.concatenate((xxx.reshape(xxx.size,1),
                                     yyy.reshape(yyy.size,1)),
                                    axis = 1)
                XX_ = poly.fit_transform(XX)
                zz = lg.predict(XX_).reshape(200,200)
                ax.plot_surface(xxx, yyy, zz, cmap=cm.coolwarm, alpha = 1) #surface plot
                
                ax.set_xlabel(feature[0])
                ax.set_ylabel(feature[1])
                ax.set_zlabel(y_label)

                figname = "%s_poly%d_%s.png" %(excelSheet,nPoly,str(feature))
                fig.savefig(figname, dpi=200) #Save plot
            else:
                figname = "none"
                
            #%% Record data in presentation
            Add_slide(nPoly,str(feature),figname,coef,"%.3f" %lg.score(X_,y),excelSheet)
            record = record.append({'num_feature':num_feature, 'feature':str(feature),'degree':nPoly,'Rsquare':np.round(lg.score(X_,y),3)}, ignore_index=True)
    
        #%% Logarithm, Exponential and Power Fitting
        if num_feature == 1:
            
            # Natural log fitting: y = a+bln(X)
            try:
                b,a = np.polyfit(np.log(xg), y, 1)
                nPoly = 0
                figname = "%s_ln_%s.png" %(excelSheet,str(feature[0]))
                r2 = R2_plot_2d(xg,y,nPoly,b,a,figname,feature[0],y_label)
                Add_slide(nPoly,feature,figname,Log_equation(feature[0],a,b),r2,excelSheet)
                record = record.append({'num_feature':num_feature, 'feature':str(feature),'degree':nPoly,'Rsquare':r2}, ignore_index=True)
            except:
                print("Contain error in log fitting")
                
            #Exponential fitting: y = a*exp(b*X)
            try:
                d,e = np.polyfit(xg, np.log(y), 1)
                nPoly = -1
                figname = "%s_exp_%s.png" %(excelSheet,str(feature[0]))
                r2 = R2_plot_2d(xg,y,nPoly,d,e,figname,feature[0],y_label)
                Add_slide(nPoly,feature,figname,Exp_equation(feature[0],d,e),r2,excelSheet)
                record = record.append({'num_feature':num_feature, 'feature':str(feature),'degree':nPoly,'Rsquare':r2}, ignore_index=True)
            except:
                print("Contain error in exponential fitting")
                
            #Power fitting: y = aX^b
            try:
                f,g = np.polyfit(np.log(xg),np.log(y),1)
                nPoly = -2
                figname = "%s_power_%s.png" %(excelSheet,str(feature[0]))
                r2 = R2_plot_2d(xg,y,nPoly,f,g,figname,feature[0],y_label)
                Add_slide(nPoly,feature,figname,Power_equation(feature[0],f,g),r2,excelSheet)
                record = record.append({'num_feature':num_feature, 'feature':str(feature),'degree':nPoly,'Rsquare':r2}, ignore_index=True)
            except:
                print("Contain error in power fitting")
                
#%% Create Summary table
record2 = record.copy()
record2['degree'] = record2['degree'].apply(lambda x:'linear' if x==1 else('poly2' if x==2 else('poly3' if x==3 else('log' if x==0 else('expo' if x==-1 else 'power')))))
nsplit = np.ceil(record.shape[0]/20)
record2 = np.array_split(record2, nsplit ,axis=0) #Split into many tables if number of row is more than 20

for rec in record2:
    fig2 = plt.figure(figsize=(8,4.5))
    ax = fig2.add_subplot(111, frame_on=False) # no visible frame
    ax.xaxis.set_visible(False)  # hide the x axis
    ax.yaxis.set_visible(False)  # hide the y axis
    tb = table(ax, rec, rowLabels=['']*rec.shape[0], loc='center')
    fig2.savefig('mytable.png',dpi = 200)
    #add summary table into presentation
    sld3 = prs.slides.add_slide(prs.slide_layouts[5])  
    sld3.placeholders[12].text = "R-squared Summary" #Add title
    graph = sld3.shapes.add_picture("mytable.png",left = Inches(1),top = Inches(1.25)) #Add plot
    graph.left = int((prs.slide_width - graph.width) / 2)
    os.remove("mytable.png") #delete image file

#%% Plot summary graph from record data
record = record[record["Rsquare"]!="Error"]
record = record.reset_index(drop=True)
record[["num_feature", "degree","Rsquare"]] = record[["num_feature", "degree","Rsquare"]].apply(pd.to_numeric) # change data type to integer
fig = plt.figure(figsize=(8,5))
ax = fig.add_subplot(111)
ax.scatter(record['degree'], record['Rsquare'])
for i, txt in enumerate(record['feature']): # label feature on each point
    ax.annotate(txt, (record['degree'][i], record['Rsquare'][i]))
ax.set_xticks(np.arange(-2,4))
ax.set_xticklabels(["Power","Expo","Log","Linear","Poly2","Poly3"])
ax.set(title = "Equation type vs R-squared plot",ylabel = "R-squared")
#fig.tight_layout()
fig.savefig("rs.png",dpi = 200)

#Add summary graph and table into presentation
sld2 = prs.slides.add_slide(prs.slide_layouts[5])  
sld2.placeholders[12].text = "R-squared Summary" #Add title
graph = sld2.shapes.add_picture("rs.png",left = Inches(1),top = Inches(1.55)) #Add plot
graph.left = int((prs.slide_width - graph.width) / 2)


os.remove("rs.png") #delete image file

prs.save('Regression_result.pptx') #save to pptx presentation

print("Finish!")

